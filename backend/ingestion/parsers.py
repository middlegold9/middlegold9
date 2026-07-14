"""把 PDF/Word/PPT 抽成纯文本（解决「数据散落」第一步）。

依赖：pdfplumber / python-docx / python-pptx
任一库缺失时该格式会被跳过，并打印提示，不影响其它格式。
"""
import os


def parse_pdf(path: str) -> str:
    try:
        import pdfplumber
    except ImportError:
        print("[skip] 未安装 pdfplumber，跳过 PDF：", path)
        return ""
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            texts.append(page.extract_text() or "")
            # 表格单独抽取（成交数据常以表格存在）
            for table in page.extract_tables():
                for row in table:
                    texts.append("\t".join(c or "" for c in row))
    return "\n".join(texts)


def parse_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        print("[skip] 未安装 python-docx，跳过 Word：", path)
        return ""
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("[skip] 未安装 python-pptx，跳过 PPT：", path)
        return ""
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def parse_file(path: str) -> str:
    """按扩展名分发解析。"""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(path)
    if ext in (".docx", ".doc"):
        return parse_docx(path)
    if ext in (".pptx", ".ppt"):
        return parse_pptx(path)
    return ""
